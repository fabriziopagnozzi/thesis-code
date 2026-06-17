from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT_PATH = (
    Path(__file__).resolve().parent
    / "@pdfs"
    / "synthetic-rag-benchmark-pipeline-reduced-slides.pptx"
)


BG = RGBColor(23, 28, 39)
PANEL = RGBColor(30, 36, 50)
WHITE = RGBColor(241, 243, 248)
MUTED = RGBColor(171, 178, 196)
BLUE = RGBColor(92, 144, 247)
BLUE_2 = RGBColor(70, 127, 244)
GREEN = RGBColor(77, 211, 160)
ORANGE = RGBColor(255, 176, 66)
LEFT_STRIP = RGBColor(78, 131, 245)

FONT_HEAD = "Aptos Display"
FONT_BODY = "Aptos"
FONT_ITALIC = "Aptos"


def set_bg(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_left_strip(slide) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(0.12), Inches(7.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = LEFT_STRIP
    shape.line.fill.background()


def add_slide_num(slide, n: int) -> None:
    box = slide.shapes.add_textbox(Inches(12.75), Inches(7.02), Inches(0.4), Inches(0.22))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = str(n)
    run.font.name = FONT_BODY
    run.font.size = Pt(12)
    run.font.color.rgb = MUTED


def add_title(slide, title: str, *, y: float = 0.62, rule_y: float = 1.36) -> None:
    font_size = 38
    if len(title) > 42:
        font_size = 33
    if len(title) > 54:
        font_size = 29

    box = slide.shapes.add_textbox(Inches(0.9), Inches(y), Inches(11.55), Inches(1.0))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = FONT_HEAD
    run.font.bold = True
    run.font.size = Pt(font_size)
    run.font.color.rgb = WHITE

    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.9), Inches(rule_y), Inches(2.55), Inches(0.045)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.fill.background()


def add_title_slide(slide, title: str, subtitle: str) -> None:
    set_bg(slide)
    add_left_strip(slide)

    box = slide.shapes.add_textbox(Inches(1.18), Inches(2.05), Inches(8.8), Inches(1.9))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = title
    r1.font.name = FONT_HEAD
    r1.font.bold = True
    r1.font.size = Pt(42)
    r1.font.color.rgb = WHITE

    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(1.08), Inches(3.9), Inches(3.05), Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.fill.background()

    sub = slide.shapes.add_textbox(Inches(1.18), Inches(4.35), Inches(5.6), Inches(0.55))
    sub.width = Inches(8.6)
    tf2 = sub.text_frame
    tf2.word_wrap = True
    tf2.clear()
    p2 = tf2.paragraphs[0]
    r2 = p2.add_run()
    r2.text = subtitle
    r2.font.name = FONT_BODY
    r2.font.size = Pt(22)
    r2.font.color.rgb = MUTED

    add_circle(slide, 9.05, 1.4, 2.35, RGBColor(42, 60, 106))
    add_circle(slide, 9.9, 4.95, 1.52, RGBColor(34, 52, 98))
    add_circle(slide, 11.55, 3.85, 1.0, RGBColor(41, 59, 105))


def add_circle(slide, x: float, y: float, size: float, color: RGBColor) -> None:
    circ = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(size), Inches(size)
    )
    circ.fill.solid()
    circ.fill.fore_color.rgb = color
    circ.line.fill.background()


def add_bullet_line(
    slide,
    bullet: str,
    text: str,
    *,
    x: float,
    y: float,
    bullet_color: RGBColor = BLUE,
    text_color: RGBColor = WHITE,
    text_size: int = 24,
    bullet_size: int = 24,
    width: float = 10.8,
    height: float = 0.72,
    italic: bool = False,
    bold_prefix: str | None = None,
) -> None:
    bullet_box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(0.35), Inches(0.28))
    btf = bullet_box.text_frame
    btf.clear()
    bp = btf.paragraphs[0]
    br = bp.add_run()
    br.text = bullet
    br.font.name = FONT_BODY
    br.font.size = Pt(bullet_size)
    br.font.color.rgb = bullet_color

    text_box = slide.shapes.add_textbox(
        Inches(x + 0.32), Inches(y - 0.01), Inches(width), Inches(height)
    )
    tf = text_box.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.space_after = Pt(0)
    if bold_prefix and text.startswith(bold_prefix):
        r1 = p.add_run()
        r1.text = bold_prefix
        r1.font.name = FONT_BODY
        r1.font.bold = True
        r1.font.size = Pt(text_size)
        r1.font.color.rgb = text_color
        r2 = p.add_run()
        r2.text = text[len(bold_prefix) :]
        r2.font.name = FONT_BODY
        r2.font.size = Pt(text_size)
        r2.font.color.rgb = text_color
        r2.font.italic = italic
    else:
        r = p.add_run()
        r.text = text
        r.font.name = FONT_BODY
        r.font.size = Pt(text_size)
        r.font.color.rgb = text_color
        r.font.italic = italic


def add_section_label(slide, text: str, *, x: float, y: float, color: RGBColor) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(4.5), Inches(0.42))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.name = FONT_BODY
    r.font.size = Pt(24)
    r.font.color.rgb = color


def add_text_panel(
    slide,
    *,
    x: float,
    y: float,
    w: float,
    h: float,
    title: str,
    lines: list[str],
    accent: RGBColor = BLUE,
    title_size: int = 18,
    body_size: int = 16,
    body_color: RGBColor = MUTED,
    body_font: str = FONT_BODY,
    bullet_lines: bool = False,
) -> None:
    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = PANEL
    panel.line.fill.background()

    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(0.1), Inches(h)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(x + 0.18), Inches(y + 0.12), Inches(w - 0.32), Inches(0.3))
    tf_title = title_box.text_frame
    tf_title.clear()
    p_title = tf_title.paragraphs[0]
    r_title = p_title.add_run()
    r_title.text = title
    r_title.font.name = FONT_BODY
    r_title.font.bold = True
    r_title.font.size = Pt(title_size)
    r_title.font.color.rgb = WHITE

    body_box = slide.shapes.add_textbox(Inches(x + 0.18), Inches(y + 0.48), Inches(w - 0.34), Inches(h - 0.58))
    tf_body = body_box.text_frame
    tf_body.word_wrap = True
    tf_body.clear()
    for idx, line in enumerate(lines):
        p = tf_body.paragraphs[0] if idx == 0 else tf_body.add_paragraph()
        p.space_after = Pt(2)
        r = p.add_run()
        r.text = f"▪ {line}" if bullet_lines else line
        r.font.name = body_font
        r.font.size = Pt(body_size)
        r.font.color.rgb = body_color


def add_two_col_artifact_slide(slide) -> None:
    set_bg(slide)
    add_title(slide, "Pipeline Overview", y=0.6, rule_y=1.36)

    blocks = [
        (
            "1. Ontology + Plans",
            [
                "LLM-generated medical ontology: conditions, subgroups, clinical axes, and value bins.",
                "Plans enumerate 2 subgroups x 2 axes, giving 4 hidden answer facets per query.",
            ],
            BLUE,
        ),
        (
            "2. Facts + Chunks",
            [
                "LLM-generated templates_data drives deterministic chunk building from structured facts.",
                "The corpus includes gold facet chunks plus explicit hard distractors and background outliers.",
            ],
            BLUE_2,
        ),
        (
            "3. Queries + Qrels",
            [
                "LLM-generated query/answer templates render the user question and canonical answer.",
                "Qrels are projected from chunk memberships, facet IDs, and cluster roles.",
            ],
            ORANGE,
        ),
        (
            "4. Embed + Evaluate",
            [
                "Embed queries and chunks, then filter for visible facets, dominant clusters, and distractors.",
                "Compare top-k, MMR, and Facility-location with aspect coverage and retrieval metrics.",
            ],
            GREEN,
        ),
    ]

    positions = [(0.98, 1.82), (6.48, 1.82), (0.98, 4.35), (6.48, 4.35)]
    for (title, lines, accent), (x, y) in zip(blocks, positions, strict=True):
        add_text_panel(
            slide,
            x=x,
            y=y,
            w=5.15,
            h=2.05,
            title=title,
            lines=lines,
            accent=accent,
            title_size=21,
            body_size=15,
            body_color=WHITE,
            bullet_lines=True,
        )


def add_stage_group_slide(
    slide,
    title: str,
    blocks: list[tuple[str, list[str], RGBColor]],
) -> None:
    set_bg(slide)
    add_title(slide, title, y=0.62, rule_y=1.37)
    top = 1.8
    block_h = 1.62
    gap = 0.2
    for idx, (head, lines, accent) in enumerate(blocks):
        y = top + idx * (block_h + gap)
        add_text_panel(
            slide,
            x=0.95,
            y=y,
            w=11.15,
            h=block_h,
            title=head,
            lines=lines,
            accent=accent,
            title_size=20,
            body_size=16,
            bullet_lines=True,
        )


def add_failure_slide(slide) -> None:
    set_bg(slide)
    add_title(slide, "Why the MIMIC Pipeline Failed", y=0.62, rule_y=1.37)

    add_text_panel(
        slide,
        x=0.98,
        y=1.82,
        w=5.35,
        h=2.2,
        title="Observed behavior",
        lines=[
            "Top-k, MMR, and Facility-location were too close to separate cleanly.",
            "Facility-location only slightly beat top-k, while MMR was slightly worse.",
            "Precision, recall, and aspect coverage all stayed weak.",
        ],
        accent=BLUE,
        title_size=20,
        body_size=17,
        bullet_lines=True,
    )
    add_text_panel(
        slide,
        x=6.62,
        y=1.82,
        w=5.55,
        h=2.2,
        title="Why coverage never emerged",
        lines=[
            "Raw discharge-note chunks were noisy and highly patient-specific.",
            "Relevant evidence was sparse and inconsistent across the candidate pool.",
            "Retrieved neighborhoods rarely exposed dominant plus complementary answer facets together.",
        ],
        accent=ORANGE,
        title_size=18,
        body_size=16,
        bullet_lines=True,
    )
    add_text_panel(
        slide,
        x=0.98,
        y=4.45,
        w=11.2,
        h=1.75,
        title="Example failure mode",
        lines=[
            "A query could retrieve several elderly-treatment chunks from one local semantic neighborhood while missing the diabetes and rehabilitation facets entirely.",
            "That makes redundancy real, but not the kind of planned multi-aspect redundancy that Facility-location is designed to exploit.",
        ],
        accent=GREEN,
        title_size=20,
        body_size=18,
        bullet_lines=True,
    )


def add_new_method_slide(slide) -> None:
    set_bg(slide)
    add_title(slide, "New Method: Curated Benchmark", y=0.62, rule_y=1.37)

    add_bullet_line(
        slide,
        "▸",
        "Start from hidden structure first, then render text",
        x=1.02,
        y=1.9,
        text_size=24,
        height=0.5,
    )
    add_bullet_line(
        slide,
        "▸",
        "Each query is built from 4 explicit answer facets: 2 subgroups x 2 clinical axes",
        x=1.02,
        y=2.38,
        text_size=24,
        height=0.76,
    )
    add_bullet_line(
        slide,
        "▸",
        "One facet is intentionally dominant; the others remain complementary; distractors are generated on purpose",
        x=1.02,
        y=3.0,
        text_size=22,
        width=10.8,
        height=0.86,
    )
    add_bullet_line(
        slide,
        "▸",
        "Queries, gold answers, and qrels all come from the same hidden plan-to-fact mapping",
        x=1.02,
        y=3.72,
        text_size=22,
        width=10.8,
        height=0.72,
    )

    add_section_label(slide, "Representative query", x=1.02, y=4.62, color=GREEN)
    add_bullet_line(
        slide,
        "▪",
        '"For patients diagnosed with encephalitis or myelitis, how do treatment duration and rehabilitation outcome differ between patients with uncomplicated diabetes and patients older than 75?"',
        x=1.45,
        y=5.08,
        bullet_color=MUTED,
        text_color=MUTED,
        text_size=18,
        bullet_size=18,
        width=10.0,
        height=0.92,
        italic=True,
    )
    add_text_panel(
        slide,
        x=1.0,
        y=6.12,
        w=11.2,
        h=0.78,
        title="Hidden facets for this query",
        lines=[
            "older than 75 x duration | older than 75 x rehab | diabetes x duration | diabetes x rehab"
        ],
        accent=GREEN,
        title_size=17,
        body_size=14,
        body_color=WHITE,
    )


def add_concrete_example_slide(slide) -> None:
    set_bg(slide)
    add_title(slide, "Concrete Example: One Query", y=0.62, rule_y=1.37)

    add_text_panel(
        slide,
        x=0.98,
        y=1.72,
        w=11.2,
        h=1.05,
        title="Query",
        lines=[
            "Among patients with encephalitis or myelitis, how do patients older than 75 and patients with uncomplicated diabetes differ in treatment duration and rehabilitation outcome?"
        ],
        accent=BLUE,
        title_size=18,
        body_size=15,
        body_color=WHITE,
    )

    add_text_panel(
        slide,
        x=0.98,
        y=2.88,
        w=3.55,
        h=1.72,
        title="chunk_0000011 | facet 1",
        lines=[
            "The 85-year-old man older than 75 presented with encephalitis or myelitis, including fever, confusion, and gait change. For total treatment duration, clinicians completed 5 days of acyclovir, while the care team documented improving mentation and reduced headache after active therapy. The discharge medication list matched the completed neurologic treatment plan and outpatient neurology follow-up."
        ],
        accent=BLUE_2,
        title_size=12,
        body_size=7,
        body_color=WHITE,
    )

    add_text_panel(
        slide,
        x=4.78,
        y=2.88,
        w=3.55,
        h=1.72,
        title="chunk_0001505 | facet 2",
        lines=[
            "The 79-year-old man older than 75 was managed for encephalitis or myelitis, with headache, altered mental status, and lower-extremity weakness. At discharge, focal weakness and impaired balance remained prominent; the discharge summary described the rehabilitation outcome as ongoing weakness and impaired balance at discharge. Neurology follow-up was arranged because recovery remained incomplete."
        ],
        accent=GREEN,
        title_size=12,
        body_size=7,
        body_color=WHITE,
    )

    add_text_panel(
        slide,
        x=8.58,
        y=2.88,
        w=3.55,
        h=1.72,
        title="chunk_0001528 | facet 3",
        lines=[
            "The 52-year-old woman with diabetes without documented end-organ complications presented with encephalitis or myelitis, including fever, confusion, and gait change. For the active therapy-course interval, acyclovir was continued for 23 days, while the treatment-duration note described improving mentation and reduced headache after active therapy. The discharge medication list matched the completed neurologic treatment plan and outpatient neurology follow-up."
        ],
        accent=BLUE_2,
        title_size=12,
        body_size=7,
        body_color=WHITE,
    )

    add_text_panel(
        slide,
        x=0.98,
        y=4.9,
        w=5.45,
        h=1.72,
        title="chunk_0001539 | facet 4",
        lines=[
            "The 69-year-old woman with uncomplicated type 2 diabetes was managed for encephalitis or myelitis, with fever, confusion, and gait change. On the day of discharge, orientation improved and gait was safe with supervised exercises; the discharge record described the functional recovery status as discharged home with outpatient neurorehabilitation. The plan emphasized caregiver teaching, home safety, and close outpatient reassessment."
        ],
        accent=BLUE_2,
        title_size=12,
        body_size=8,
        body_color=WHITE,
    )

    add_text_panel(
        slide,
        x=6.68,
        y=4.9,
        w=5.45,
        h=1.72,
        title="chunk_0001571 | distractor",
        lines=[
            "The 52-year-old man with baseline dementia with memory impairment presented with pneumonia, including dyspnea, leukocytosis, and radiographic consolidation. For the active therapy-course interval, records showed 7 days of treatment with azithromycin, while the treatment-duration note described improving dyspnea and downtrending fever after active therapy. Respiratory status was stable, and no additional inpatient antibiotic escalation was documented."
        ],
        accent=ORANGE,
        title_size=12,
        body_size=8,
        body_color=WHITE,
    )


def add_expected_behavior_slide(slide) -> None:
    set_bg(slide)
    add_title(slide, "Designed Behavior", y=0.62, rule_y=1.37)

    add_bullet_line(
        slide,
        "▸",
        "top-k can be relevant but redundant",
        x=1.02,
        y=1.9,
        bullet_color=BLUE,
        text_size=25,
        bold_prefix="top-k",
    )
    add_bullet_line(
        slide,
        "▪",
        "It often fills the context budget with repeated chunks from the dominant gold facet",
        x=1.45,
        y=2.34,
        bullet_color=MUTED,
        text_color=MUTED,
        text_size=20,
        bullet_size=18,
        width=9.9,
    )

    add_bullet_line(
        slide,
        "▸",
        "MMR can reduce redundancy but drift toward distractors",
        x=1.02,
        y=3.02,
        bullet_color=ORANGE,
        text_size=25,
        bold_prefix="MMR",
    )
    add_bullet_line(
        slide,
        "▪",
        "Dispersion is not the same as covering the remaining answer facets",
        x=1.45,
        y=3.46,
        bullet_color=MUTED,
        text_color=MUTED,
        text_size=20,
        bullet_size=18,
        width=9.9,
    )

    add_bullet_line(
        slide,
        "▸",
        "Facility-location should cover the remaining local clusters",
        x=1.02,
        y=4.14,
        bullet_color=GREEN,
        text_size=25,
        bold_prefix="Facility-location",
    )
    add_bullet_line(
        slide,
        "▪",
        "The new pipeline makes those dominant, complementary, and distractor clusters explicit per query.",
        x=1.45,
        y=4.58,
        bullet_color=MUTED,
        text_color=MUTED,
        text_size=20,
        bullet_size=18,
        width=9.8,
    )


def build_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    add_title_slide(
        slide,
        "Synthetic Clinical RAG Benchmark",
        "From failed MIMIC-IV notes to a coverage-sensitive curated pipeline",
    )
    add_slide_num(slide, 1)

    slide = prs.slides.add_slide(blank)
    add_failure_slide(slide)
    add_slide_num(slide, 2)

    slide = prs.slides.add_slide(blank)
    add_new_method_slide(slide)
    add_slide_num(slide, 3)

    slide = prs.slides.add_slide(blank)
    add_concrete_example_slide(slide)
    add_slide_num(slide, 4)

    slide = prs.slides.add_slide(blank)
    add_two_col_artifact_slide(slide)
    add_slide_num(slide, 5)

    slide = prs.slides.add_slide(blank)
    add_stage_group_slide(
        slide,
        "Stages 1-3: Retrieval Geometry",
        [
            (
                "Stage 1. Ontology + configuration",
                [
                    "Ontology defines conditions, subgroup rules, and exactly two axes: treatment duration and rehab outcome.",
                    "Example condition entry: encephalitis or myelitis with acyclovir / corticosteroids and short / standard / prolonged bins.",
                ],
                BLUE,
            ),
            (
                "Stage 2. Hidden query plans",
                [
                    "Enumerate condition and subgroup-pair combinations; every query gets 4 facets.",
                    "Example plan: q00001 compares age_over_75 versus uncomplicated_diabetes for encephalitis_myelitis.",
                ],
                BLUE_2,
            ),
            (
                "Stage 3. Dominance calibration",
                [
                    "Optional probe embeddings pick the facet that is naturally closest to the query.",
                    "Dominance becomes an embedding-space property, not a wording trick in the final query text.",
                ],
                GREEN,
            ),
        ],
    )
    add_slide_num(slide, 6)

    slide = prs.slides.add_slide(blank)
    add_stage_group_slide(
        slide,
        "Stages 4-7: Evidence And Labels",
        [
            (
                "Stage 4. Structured clinical facts",
                [
                    "Expand each facet into gold fact rows; add hard negatives and background outlier clusters.",
                    "Example fact: q00001_f3 -> prolonged, 24 days, acyclovir, complementary_gold.",
                ],
                BLUE,
            ),
            (
                "Stage 5. Chunk rendering + normalization",
                [
                    "Render validated chunk documents with deterministic templates, direct LLM generation, or LLM rewrites.",
                    "Example chunk mentions encephalitis or myelitis, subgroup evidence, and the duration value required by validation.",
                ],
                ORANGE,
            ),
            (
                "Stage 6-7. Queries, answers, qrels",
                [
                    "Render the user query and canonical answer from query_answer_templates.yaml tied to the same hidden plan.",
                    "Project qrels directly from chunk memberships instead of inferring them from text after the fact.",
                ],
                GREEN,
            ),
        ],
    )
    add_slide_num(slide, 7)

    slide = prs.slides.add_slide(blank)
    add_stage_group_slide(
        slide,
        "Stages 8-11: Validation",
        [
            (
                "Stage 8. Embeddings",
                [
                    "Embed unique chunk documents and queries with the configured sentence-embedding model.",
                    "All later retrieval scores operate on the same candidate space.",
                ],
                BLUE,
            ),
            (
                "Stage 9-10. Geometry filter + evaluation",
                [
                    "Keep only queries where all facets are visible, top-k over-selects the dominant facet, and distractors are present.",
                    "Evaluate top-k, MMR, and Facility-location over the same top-N semantic pool.",
                ],
                BLUE_2,
            ),
            (
                "Stage 11. Diagnostics + plots",
                [
                    "Embedding-geometry plots explain local cluster structure and dominant-facet pressure.",
                    "The benchmark is only useful if the measured behavior is coverage-sensitive.",
                ],
                GREEN,
            ),
        ],
    )
    add_slide_num(slide, 8)

    slide = prs.slides.add_slide(blank)
    add_expected_behavior_slide(slide)
    add_slide_num(slide, 9)

    return prs


def main() -> None:
    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs = build_deck()
    prs.save(OUT_PATH)
    print(f"[write] {OUT_PATH}")


if __name__ == "__main__":
    main()
